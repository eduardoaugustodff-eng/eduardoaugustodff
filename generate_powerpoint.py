from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from data import catalogData

# Cores da empresa
DARK_BG = RGBColor(10, 14, 39)        # #0a0e27
PRIMARY_DARK = RGBColor(0, 31, 63)    # #001f3f
PRIMARY_BLUE = RGBColor(0, 102, 204)  # #0066cc
ACCENT_BLUE = RGBColor(0, 153, 255)  # #0099ff
TEXT_PRIMARY = RGBColor(255, 255, 255)
TEXT_SECONDARY = RGBColor(179, 179, 204)
SUCCESS_GREEN = RGBColor(0, 212, 170)

def create_catalog_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Capa
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_DARK

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "AUTO PEÇAS NUNEDIESEL"
    p.font.size = Pt(66)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "Catálogo de Peças - Motor Cummins 6BT"
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT_SECONDARY
    p.alignment = PP_ALIGN.CENTER

    # Descrição
    desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.2), Inches(9), Inches(1))
    desc_frame = desc_box.text_frame
    p = desc_frame.paragraphs[0]
    p.text = "Ferramenta para Vendedores"
    p.font.size = Pt(18)
    p.font.color.rgb = SUCCESS_GREEN
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: Índice
    slide = add_category_slide(prs, "Índice de Categorias")

    items_text = ""
    for i, category in enumerate(catalogData, 1):
        items_text += f"{i}. {category['icon']} {category['category']} ({len(category['items'])} peças)\n"

    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, category in enumerate(catalogData, 1):
        if i > 1:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i-1]
        p.text = f"{i}. {category['icon']} {category['category']}"
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = ACCENT_BLUE
        p.space_after = Pt(8)
        p.font.bold = True

    # Slides das categorias
    for category in catalogData:
        create_category_slide(prs, category)

    # Slide final: Contato
    slide = add_category_slide(prs, "Informações")

    info_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(7), Inches(4))
    text_frame = info_box.text_frame
    text_frame.word_wrap = True

    info_items = [
        "📞 CONTATO",
        "",
        "Auto Peças Nunediesel",
        "",
        "💬 Entre em contato para:",
        "  • Consultar disponibilidade",
        "  • Solicitar orçamento",
        "  • Agendar entrega",
        "",
        "Obrigado pela confiança! ✅"
    ]

    for i, item in enumerate(info_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18 if "CONTATO" in item or "Obrigado" in item else 14)
        p.font.color.rgb = ACCENT_BLUE if "CONTATO" in item or "Obrigado" in item else TEXT_PRIMARY
        p.font.bold = "CONTATO" in item or "Obrigado" in item
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)

    # Salvar apresentação
    prs.save('Catalogo_Nunediesel_Cummins_6BT.pptx')
    print("✅ Apresentação criada: Catalogo_Nunediesel_Cummins_6BT.pptx")

def add_category_slide(prs, title):
    """Cria um slide com cabeçalho padrão"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

    # Barra superior
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_DARK
    shape.line.color.rgb = ACCENT_BLUE
    shape.line.width = Pt(2)

    # Título
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE

    return slide

def create_category_slide(prs, category):
    """Cria um slide para cada categoria"""
    slide = add_category_slide(prs, f"{category['icon']} {category['category']}")

    # Contador de peças
    count_box = slide.shapes.add_textbox(Inches(8.5), Inches(0.2), Inches(1.3), Inches(0.5))
    count_frame = count_box.text_frame
    p = count_frame.paragraphs[0]
    p.text = f"{len(category['items'])} peças"
    p.font.size = Pt(12)
    p.font.color.rgb = SUCCESS_GREEN
    p.alignment = PP_ALIGN.RIGHT

    # Conteúdo
    y_position = 1.2
    row_height = 0.45

    # Cabeçalho da tabela
    header_y = y_position
    headers = ["Código", "Descrição", "Nº Peça"]
    header_widths = [2.2, 4.5, 2.8]
    x_pos = 0.3

    for header, width in zip(headers, header_widths):
        header_box = slide.shapes.add_textbox(Inches(x_pos), Inches(header_y), Inches(width), Inches(0.35))
        header_frame = header_box.text_frame
        header_frame.word_wrap = False
        p = header_frame.paragraphs[0]
        p.text = header
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = ACCENT_BLUE
        x_pos += width

    # Linha separadora
    line_y = header_y + 0.35
    line = slide.shapes.add_connector(1, Inches(0.3), Inches(line_y), Inches(9.7), Inches(line_y))
    line.line.color.rgb = ACCENT_BLUE
    line.line.width = Pt(1)

    # Items
    y_position = line_y + 0.15
    for item in category['items']:
        if y_position > 7:  # Se ultrapassar a página, para
            break

        x_pos = 0.3

        # Código
        code_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_position), Inches(2.2), Inches(row_height))
        code_frame = code_box.text_frame
        code_frame.word_wrap = True
        p = code_frame.paragraphs[0]
        p.text = item['code']
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = SUCCESS_GREEN

        # Descrição
        x_pos += 2.2
        desc_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_position), Inches(4.5), Inches(row_height))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        p = desc_frame.paragraphs[0]
        p.text = item['description']
        if item.get('note'):
            p.text += f" ({item['note']})"
        p.font.size = Pt(8)
        p.font.color.rgb = TEXT_PRIMARY

        # Número de peça
        x_pos += 4.5
        part_box = slide.shapes.add_textbox(Inches(x_pos), Inches(y_position), Inches(2.8), Inches(row_height))
        part_frame = part_box.text_frame
        part_frame.word_wrap = True
        p = part_frame.paragraphs[0]
        p.text = item['partNumber'] if item['partNumber'] else "-"
        p.font.size = Pt(8)
        p.font.color.rgb = TEXT_SECONDARY

        y_position += row_height + 0.05

if __name__ == "__main__":
    create_catalog_presentation()
